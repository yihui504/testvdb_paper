#!/usr/bin/env python3
"""21-page GEAR-style .pptx — P3 slim + P4/5 merge (22->21)."""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG=os.path.join(os.path.dirname(__file__),"figures")
OUT=os.path.join(os.path.dirname(__file__),"TestVDB_slides.pptx")
HBG=RGBColor(0x1B,0x3A,0x5C); AB=RGBColor(0x2E,0x5C,0x8A); AO=RGBColor(0xD9,0x8E,0x48)
AG=RGBColor(0x4A,0x8C,0x5C); AR=RGBColor(0xB0,0x50,0x50); TD=RGBColor(0x33,0x33,0x33)
TG=RGBColor(0x66,0x66,0x66); LB=RGBColor(0xF0,0xF4,0xF8); WH=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def header(sl,title,n,sec):
    bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb=HBG; bar.line.fill.background()
    tb=sl.shapes.add_textbox(Inches(0.45),Inches(0.12),Inches(10.5),Inches(0.85))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(22); p.font.bold=True; p.font.color.rgb=WH
    pb=sl.shapes.add_textbox(Inches(10.8),Inches(0.12),Inches(2.3),Inches(0.45))
    pp=pb.text_frame.paragraphs[0]; pp.text=f"P{n}/20 {sec}"; pp.alignment=PP_ALIGN.RIGHT
    pp.font.size=Pt(11); pp.font.color.rgb=RGBColor(0xBB,0xCC,0xDD)
    st=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(1.05),prs.slide_width,Inches(0.06))
    st.fill.solid(); st.fill.fore_color.rgb=AO; st.line.fill.background()

def bullets(sl,items,top=1.5,left=0.7,w=12,h=5.5,sz=16):
    tb=sl.shapes.add_textbox(Inches(left),Inches(top),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        t,c=(item if isinstance(item,tuple) else (item,TD))
        p.text="•  "+t; p.font.size=Pt(sz); p.font.color.rgb=c; p.space_after=Pt(10)

def table(sl,hs,rs,top=1.5,left=0.6,w=12.1,h=5.2):
    g=sl.shapes.add_table(len(rs)+1,len(hs),Inches(left),Inches(top),Inches(w),Inches(h)).table
    g.first_row=True; g.horz_banding=False
    for j,h in enumerate(hs):
        c=g.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
        for p in c.text_frame.paragraphs: p.font.bold=True; p.font.size=Pt(13); p.font.color.rgb=WH; p.alignment=PP_ALIGN.CENTER
    for i,row in enumerate(rs):
        for j,v in enumerate(row):
            c=g.cell(i+1,j); c.text=str(v); c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
            for p in c.text_frame.paragraphs: p.font.size=Pt(11.5); p.font.color.rgb=TD

def image(sl,name,top=1.35,left=1.4,w=10.5):
    sl.shapes.add_picture(os.path.join(FIG,name),Inches(left),Inches(top),Inches(w))

def bn(sl,num,label,sub=None,top=2.2):
    tb=sl.shapes.add_textbox(Inches(1),Inches(top),Inches(11.3),Inches(2.5))
    p=tb.text_frame.paragraphs[0]; p.text=num; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(96); p.font.bold=True; p.font.color.rgb=AO
    lb=sl.shapes.add_textbox(Inches(1),Inches(top+2.3),Inches(11.3),Inches(1))
    lp=lb.text_frame.paragraphs[0]; lp.text=label; lp.alignment=PP_ALIGN.CENTER; lp.font.size=Pt(20); lp.font.color.rgb=TD
    if sub:
        sb=sl.shapes.add_textbox(Inches(1),Inches(top+3.1),Inches(11.3),Inches(1))
        sp=sb.text_frame.paragraphs[0]; sp.text=sub; sp.alignment=PP_ALIGN.CENTER; sp.font.size=Pt(14); sp.font.color.rgb=TG

def two_col(sl,lt,li,rt,ri,lc=AR,rc=AG):
    lb=sl.shapes.add_textbox(Inches(0.6),Inches(1.4),Inches(6),Inches(0.6))
    lp=lb.text_frame.paragraphs[0]; lp.text=lt; lp.font.size=Pt(16); lp.font.bold=True; lp.font.color.rgb=lc
    bullets(sl,li,top=2.0,left=0.7,w=5.8,h=4.8,sz=14)
    dv=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(6.5),Inches(1.6),Inches(0.04),Inches(5))
    dv.fill.solid(); dv.fill.fore_color.rgb=RGBColor(0xDD,0xDD,0xDD); dv.line.fill.background()
    rb=sl.shapes.add_textbox(Inches(6.9),Inches(1.4),Inches(6),Inches(0.6))
    rp=rb.text_frame.paragraphs[0]; rp.text=rt; rp.font.size=Pt(16); rp.font.bold=True; rp.font.color.rgb=rc
    bullets(sl,ri,top=2.0,left=7.0,w=5.8,h=4.8,sz=14)

def slide(n,sec,title):
    s=prs.slides.add_slide(BLANK); header(s,title,n,sec); return s

def note(sl,text,top=6.6):
    nb=sl.shapes.add_textbox(Inches(0.6),Inches(top),Inches(12),Inches(0.5))
    p=nb.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(11); p.font.color.rgb=TG

def naive_pipeline(sl):
    boxes=[(Inches(0.5),"1. Extract","LLM reads documentation\n-> behavioral claims",AB),
           (Inches(4.7),"2. Attack","boundary inputs\ntargeting each claim",AB),
           (Inches(8.9),"3. Judge","LLM checks response\n-> conform or violate?",AO)]
    for left,title,desc,color in boxes:
        box=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left,Inches(2.5),Inches(3.2),Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb=LB; box.line.color.rgb=color; box.line.width=Pt(2)
        tf=box.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.text=title; p.font.size=Pt(18); p.font.bold=True; p.font.color.rgb=color
        p2=tf.add_paragraph(); p2.text=desc; p2.font.size=Pt(13); p2.font.color.rgb=TD
    for left in [Inches(3.9),Inches(8.1)]:
        arr=sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,left,Inches(3.3),Inches(0.6),Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb=AO; arr.line.fill.background()
    wb=sl.shapes.add_textbox(Inches(8.9),Inches(4.9),Inches(3.2),Inches(0.5))
    wp=wb.text_frame.paragraphs[0]; wp.text="-> verdict (may be wrong)"; wp.font.size=Pt(14); wp.font.color.rgb=AR; wp.font.bold=True

def defect_example(sl):
    """nprobe=0 doc-vs-actual defect illustration (native shapes)."""
    # Left box (green = documentation)
    lb=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.0),Inches(4.5),Inches(2.2))
    lb.fill.solid(); lb.fill.fore_color.rgb=LB; lb.line.color.rgb=AG; lb.line.width=Pt(2)
    tf=lb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text="Documentation"; p.font.size=Pt(14); p.font.bold=True; p.font.color.rgb=AG
    p2=tf.add_paragraph(); p2.text='"nprobe (int):\n Range 0-16384"'; p2.font.size=Pt(13); p2.font.color.rgb=TD
    p3=tf.add_paragraph(); p3.text="(should reject nprobe=0)"; p3.font.size=Pt(12); p3.font.color.rgb=TG
    # Arrow
    arr=sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(5.2),Inches(3.7),Inches(0.8),Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb=AR; arr.line.fill.background()
    # Right box (red = actual behavior)
    rb=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.2),Inches(3.0),Inches(5.5),Inches(2.2))
    rb.fill.solid(); rb.fill.fore_color.rgb=LB; rb.line.color.rgb=AR; rb.line.width=Pt(2)
    tf2=rb.text_frame; tf2.word_wrap=True
    p=tf2.paragraphs[0]; p.text="API Actual Behaviour"; p.font.size=Pt(14); p.font.bold=True; p.font.color.rgb=AR
    p2=tf2.add_paragraph(); p2.text="accepts nprobe=0"; p2.font.size=Pt(13); p2.font.color.rgb=TD
    p3=tf2.add_paragraph(); p3.text="-> search runs (wrong results)"; p3.font.size=Pt(13); p3.font.color.rgb=AR
    p4=tf2.add_paragraph(); p4.text="no crash, no error code"; p4.font.size=Pt(12); p4.font.color.rgb=TG
    # Bottom annotation
    nb=sl.shapes.add_textbox(Inches(0.5),Inches(5.5),Inches(11.5),Inches(0.6))
    np_=nb.text_frame.paragraphs[0]
    np_.text="Divergence -> query semantics corrupted; 37 of 38 such defects produce no crash"
    np_.font.size=Pt(14); np_.font.color.rgb=AR; np_.font.bold=True

# === 21 slides ===

s=prs.slides.add_slide(BLANK)
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb=HBG; bg.line.fill.background()
tb=s.shapes.add_textbox(Inches(0.8),Inches(2.4),Inches(11.7),Inches(2.2)); tf=tb.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="TestVDB"; p.font.size=Pt(54); p.font.bold=True; p.font.color.rgb=WH
p2=tf.add_paragraph(); p2.text="Using LLMs to Find Bugs Where Actual Behaviour Diverges from Its Documentation"
p2.font.size=Pt(22); p2.font.color.rgb=RGBColor(0xCC,0xDD,0xEE)
meta=s.shapes.add_textbox(Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.5))
for i,line in enumerate(["Authors (TBD)","Affiliation","Venue / Session (TBD)"]):
    p=meta.text_frame.paragraphs[0] if i==0 else meta.text_frame.add_paragraph()
    p.text=line; p.font.size=Pt(16); p.font.color.rgb=RGBColor(0xAA,0xCC,0xDD)

s=slide(2,"Opening","VDBMS defects are costly — and mostly functional")
image(s,"fig1_rag_arch.png",top=1.3,left=2.5,w=8.3)
table(s,["Source","Finding"],[["Empirical bug study (Xie 2025)","> 50% functional failures"],["Testing roadmap (Wang 2025)","~43% incorrect behavior; oracle key challenge"],["VDBFuzz (Wang 2026)","only VDBMS fuzzer — crash oracle"]],top=4.7,h=2.2)

# P3 slim (2 bullets + nprobe defect example)
s=slide(3,"Problem","Doc-implementation defects: actual behaviour diverges from documentation")
bullets(s,[("Doc-implementation consistency — does API accept/reject match documentation?",TD),("Correctness — is the result mathematically right?",TD)],top=1.3,sz=16)
defect_example(s)

# P4 merged (fig2 cases + 37/38 big number)
s=slide(4,"Problem","37 of 38 defects do not crash — fuzzers miss them")
tb=s.shapes.add_textbox(Inches(1),Inches(5.0),Inches(11.3),Inches(1.2))
p=tb.text_frame.paragraphs[0]; p.text="37 / 38"; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(64); p.font.bold=True; p.font.color.rgb=AO
lb=s.shapes.add_textbox(Inches(1),Inches(6.2),Inches(11.3),Inches(0.5))
lp=lb.text_frame.paragraphs[0]; lp.text="acknowledged defects produce no crash — fuzzers cannot reach this class"
lp.alignment=PP_ALIGN.CENTER; lp.font.size=Pt(16); lp.font.color.rgb=TD

# P5 Table 1 (was P6)
# P6 naive pipeline (was P7)
s=slide(5,"Method","TestVDB: LLM extracts claims and judges conformance")
naive_pipeline(s)
note(s,"Crash/differential/metamorphic/property-based oracles all miss this residual. Only an LLM can judge accept/reject vs documentation. This naive pipeline produces false positives — the LLM judge is unreliable (next).")

# P7-P10 core insight (was P8-P11)
s=slide(6,"Core insight","The source-ambiguity gap: assertions vs claims")
image(s,"fig3_source_ambiguity_gap.png",top=1.3,left=0.8,w=11.7)

s=slide(7,"Core insight","Family-specific errors: the judge confirms the extractor's biases")
bullets(s,["one LLM family extracts claims AND judges — shared biases","judge confirms extractor errors — self-preference","Panickssery (2024); Wataoka (2024)",("mitigation: cross-model validation",AG)])

s=slide(8,"Core insight","Task-intrinsic errors: different families infer the same wrong claim")
image(s,"fig4_two_layer_venn.png",top=1.3,left=0.9,w=11.5)

s=slide(9,"Core insight","Cross-model validation covers family-specific, not task-intrinsic")
two_col(s,"family-specific",["consistencyLevel — GLM/DeepSeek disagree","cross-model catches divergence"],
     "task-intrinsic",["timeout — both extract '>= 1' (same error)","cross-model sees agreement","only source falsifies"],lc=AB,rc=AO)

# P11 complete pipeline architecture (was P12)
s=slide(10,"Method","Complete pipeline: source-grounded falsification resolves what cross-model cannot")
image(s,"fig5_pipeline.png",top=1.3,left=0.5,w=12.3)
note(s,"Added to naive pipeline (P6): dev-reviewer reads source to falsify LLM claims. 12-clause pilot: cross-model 7/12, source 12/12 (see P17 for n=29).")

# P12 falsification rule (was P13)
s=slide(11,"Method","Falsification rule: shardsNum=0 selects the default")
two_col(s,"Over-strict clause (LLM)",["shardsNum >= 1","probe: shardsNum=0 -> API 200","LLM verdict: 'violation'"],
     "Source-grounded falsification",["source: if shardsNum==0 { use default }","0 selects default — over-strict","FP killed","opposite of MASTOR (as currently designed)"],lc=AR,rc=AG)
note(s,"+ dev-reviewer reads source to falsify claims; novelty gate removes duplicates. Suppresses FPs from naive pipeline (P6).")

# P13-P21 evaluation + closing (was P14-P22, all -1)
s=slide(12,"Evaluation","Four RQs; 111 submitted, 38 acknowledged")
g1=s.shapes.add_table(5,2,Inches(0.4),Inches(1.5),Inches(5.8),Inches(3.8)).table; g1.first_row=True; g1.horz_banding=False
for j,h in enumerate(["RQ","Question"]):
    c=g1.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
    for p in c.text_frame.paragraphs: p.font.bold=True; p.font.size=Pt(12); p.font.color.rgb=WH
for i,(rq,q) in enumerate([["RQ1","how many doc-impl defects?"],["RQ2","does source suppress FPs?"],["RQ3","cross-model vs source?"],["RQ4","model-free subclass?"]]):
    for j,v in enumerate([rq,q]):
        c=g1.cell(i+1,j); c.text=v; c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
        for p in c.text_frame.paragraphs: p.font.size=Pt(11); p.font.color.rgb=TD
g2=s.shapes.add_table(7,3,Inches(6.8),Inches(1.5),Inches(6.1),Inches(4.5)).table; g2.first_row=True; g2.horz_banding=False
for j,h in enumerate(["VDBMS","Submitted","Ack."]):
    c=g2.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
    for p in c.text_frame.paragraphs: p.font.bold=True; p.font.size=Pt(12); p.font.color.rgb=WH
for i,row in enumerate([["Milvus","51","22"],["Weaviate","30","3"],["Qdrant","26","13"],["MeiliSearch","3","0"],["Chroma","1","0"],["Total","111","38"]]):
    for j,v in enumerate(row):
        c=g2.cell(i+1,j); c.text=v; c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
        for p in c.text_frame.paragraphs: p.font.size=Pt(11); p.font.color.rgb=TD

s=slide(13,"Evaluation","~85% doc-impl defects; VDBFuzz: 0 crashes")
two_col(s,"Composition (not prevalence)",["~85% doc-implementation defects","~10% classical","~5% concurrency","89% on acknowledged"],
     "VDBFuzz (Qdrant v1.18.2)",["we ran VDBFuzz: 26,000 requests","0 crashes, 0 non-200","TestVDB found doc-impl defects","disjoint classes"],lc=AB,rc=AG)

s=slide(14,"Evaluation","Source anchor: 81% FP suppression (up from 31%)")
bn(s,"81%","FP suppressed by source anchor","up from 31%; 96.7% TP retention (n=30)")

s=slide(15,"Evaluation","Precision scales: 25.5% -> 45.6% -> 69.2%")
image(s,"fig7_ablation.png",top=1.5,left=2.5,w=8.3)

s=slide(16,"Evaluation","RQ3 n=29: source 16/16; explicit-bound 0/13; kappa=1.0")
table(s,["Subtype","n","Task-intrinsic","Cross-model","Source"],
      [["Parameter over-strict","12","5/12","7/12","12/12"],["Behavior over-strict","4","4/4","0/4","4/4"],
       ["Explicit-bound negative","13","0/13","-","-"],["Within-vendor contrast","-","56% vs 0%","-","-"],
       ["Cross-model kappa (n=20)","20","-","kappa=1.0","-"]],top=1.5)

s=slide(17,"Evaluation","A model-free invariant subclass finds bugs on its own")
table(s,["Invariant","Observation","Cross-vendor"],[["COSINE bound","distance > 1.0 identical vectors","Milvus+Qdrant"],
     ["Index completeness","index returns 2 of 25","Milvus+Qdrant"],["Payload filter","filter absent field returns points","Milvus+Qdrant"]])

s=slide(18,"Related work","Prior work: low-ambiguity; TestVDB: ambiguous regime")
table(s,["Line of work","Representative","Difference"],
      [["VDBMS testing","VDBFuzz/roadmap/bug study","crash vs doc-impl"],
       ["REST oracle","AGORA+/SATORI/MASTOR","low-ambiguity; we use NL docs"],
       ["Doc-derived oracle","Toradocu/Doc2OracLL/AugmenTest/Konstantinou","trust LLM; we falsify"],
       ["DB correctness","NoREC/TLP/DQE/DDLCheck","reference semantics; ours lacks"],
       ["LLM-judge","Panickssery/Wataoka/Haldar","orthogonal — source grounding"]],top=1.4,h=5.6)

s=slide(19,"Closing","Threats and conclusion")
two_col(s,"Threats",["over-strict (n=16) most contingent","Milvus+Qdrant only","mechanism correlative"],
     "Conclusion",["doc-impl resists checking","LLM: family-specific + task-intrinsic","generalizes beyond VDBMSs"],lc=AR,rc=AG)

s=slide(20,"Closing","TestVDB — four core results")
image(s,"fig8_summary.png",top=1.4,left=1.5,w=10.3)

prs.save(OUT)
print(f"[wrote] {OUT} ({len(prs.slides._sldIdLst)} slides)")
