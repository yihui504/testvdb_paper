#!/usr/bin/env python3
"""21-page Chinese GEAR-style .pptx — P3 slim + P4/5 merge. CJK: 微软雅黑."""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIG=os.path.join(os.path.dirname(__file__),"figures")
OUT=os.path.join(os.path.dirname(__file__),"TestVDB_slides_zh.pptx")
CJK="微软雅黑"
HBG=RGBColor(0x1B,0x3A,0x5C); AB=RGBColor(0x2E,0x5C,0x8A); AO=RGBColor(0xD9,0x8E,0x48)
AG=RGBColor(0x4A,0x8C,0x5C); AR=RGBColor(0xB0,0x50,0x50); TD=RGBColor(0x33,0x33,0x33)
TG=RGBColor(0x66,0x66,0x66); LB=RGBColor(0xF0,0xF4,0xF8); WH=RGBColor(0xFF,0xFF,0xFF)
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def _f(p,s,c,b=False):
    p.font.name=CJK; p.font.size=Pt(s); p.font.color.rgb=c; p.font.bold=b

def header(sl,t,n,sec):
    bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb=HBG; bar.line.fill.background()
    tb=sl.shapes.add_textbox(Inches(0.45),Inches(0.12),Inches(10.5),Inches(0.85))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=t; _f(p,20,WH,True)
    pb=sl.shapes.add_textbox(Inches(10.8),Inches(0.12),Inches(2.3),Inches(0.45))
    pp=pb.text_frame.paragraphs[0]; pp.text=f"P{n}/20 {sec}"; pp.alignment=PP_ALIGN.RIGHT; _f(pp,11,RGBColor(0xBB,0xCC,0xDD))
    st=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,Inches(1.05),prs.slide_width,Inches(0.06))
    st.fill.solid(); st.fill.fore_color.rgb=AO; st.line.fill.background()

def bullets(sl,items,top=1.5,left=0.7,w=12,h=5.5,sz=16):
    tb=sl.shapes.add_textbox(Inches(left),Inches(top),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        t,c=(item if isinstance(item,tuple) else (item,TD))
        p.text="•  "+t; _f(p,sz,c); p.space_after=Pt(10)

def table(sl,hs,rs,top=1.5,left=0.6,w=12.1,h=5.2):
    g=sl.shapes.add_table(len(rs)+1,len(hs),Inches(left),Inches(top),Inches(w),Inches(h)).table
    g.first_row=True; g.horz_banding=False
    for j,h in enumerate(hs):
        c=g.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
        for p in c.text_frame.paragraphs: _f(p,13,WH,True); p.alignment=PP_ALIGN.CENTER
    for i,row in enumerate(rs):
        for j,v in enumerate(row):
            c=g.cell(i+1,j); c.text=str(v); c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
            for p in c.text_frame.paragraphs: _f(p,11.5,TD)

def image(sl,name,top=1.35,left=1.4,w=10.5):
    sl.shapes.add_picture(os.path.join(FIG,name),Inches(left),Inches(top),Inches(w))

def bn(sl,num,label,sub=None,top=2.2):
    tb=sl.shapes.add_textbox(Inches(1),Inches(top),Inches(11.3),Inches(2.5))
    p=tb.text_frame.paragraphs[0]; p.text=num; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(96); p.font.bold=True; p.font.color.rgb=AO
    lb=sl.shapes.add_textbox(Inches(1),Inches(top+2.3),Inches(11.3),Inches(1))
    lp=lb.text_frame.paragraphs[0]; lp.text=label; lp.alignment=PP_ALIGN.CENTER; _f(lp,20,TD)
    if sub:
        sb=sl.shapes.add_textbox(Inches(1),Inches(top+3.1),Inches(11.3),Inches(1))
        sp=sb.text_frame.paragraphs[0]; sp.text=sub; sp.alignment=PP_ALIGN.CENTER; _f(sp,14,TG)

def tc(sl,lt,li,rt,ri,lc=AR,rc=AG):
    lb=sl.shapes.add_textbox(Inches(0.6),Inches(1.4),Inches(6),Inches(0.6))
    lp=lb.text_frame.paragraphs[0]; lp.text=lt; _f(lp,16,lc,True)
    bullets(sl,li,top=2.0,left=0.7,w=5.8,h=4.8,sz=14)
    dv=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(6.5),Inches(1.6),Inches(0.04),Inches(5))
    dv.fill.solid(); dv.fill.fore_color.rgb=RGBColor(0xDD,0xDD,0xDD); dv.line.fill.background()
    rb=sl.shapes.add_textbox(Inches(6.9),Inches(1.4),Inches(6),Inches(0.6))
    rp=rb.text_frame.paragraphs[0]; rp.text=rt; _f(rp,16,rc,True)
    bullets(sl,ri,top=2.0,left=7.0,w=5.8,h=4.8,sz=14)

def slide(n,sec,title):
    s=prs.slides.add_slide(BLANK); header(s,title,n,sec); return s

def note(sl,text,top=6.6):
    nb=sl.shapes.add_textbox(Inches(0.6),Inches(top),Inches(12),Inches(0.5))
    p=nb.text_frame.paragraphs[0]; p.text=text; _f(p,11,TG)

def naive_pipeline(sl):
    boxes=[(Inches(0.5),"1. 提取","LLM 读文档\n→ 行为 claim",AB),(Inches(4.7),"2. 攻击","边界输入\n针对每个 claim",AB),(Inches(8.9),"3. 裁决","LLM 检查响应\n→ 符合或违反？",AO)]
    for left,title,desc,color in boxes:
        box=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,left,Inches(2.5),Inches(3.2),Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb=LB; box.line.color.rgb=color; box.line.width=Pt(2)
        tf=box.text_frame; tf.word_wrap=True
        p=tf.paragraphs[0]; p.text=title; _f(p,18,color,True); p2=tf.add_paragraph(); p2.text=desc; _f(p2,13,TD)
    for left in [Inches(3.9),Inches(8.1)]:
        arr=sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,left,Inches(3.3),Inches(0.6),Inches(0.4))
        arr.fill.solid(); arr.fill.fore_color.rgb=AO; arr.line.fill.background()
    wb=sl.shapes.add_textbox(Inches(8.9),Inches(4.9),Inches(3.2),Inches(0.5))
    wp=wb.text_frame.paragraphs[0]; wp.text="→ 裁决（可能错误）"; _f(wp,14,AR,True)

def defect_example(sl):
    lb=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.5),Inches(3.0),Inches(4.5),Inches(2.2))
    lb.fill.solid(); lb.fill.fore_color.rgb=LB; lb.line.color.rgb=AG; lb.line.width=Pt(2)
    tf=lb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text="文档"; _f(p,14,AG,True)
    p2=tf.add_paragraph(); p2.text='"nprobe (int):\n 范围 0-16384"'; _f(p2,13,TD)
    p3=tf.add_paragraph(); p3.text="（应拒绝 nprobe=0）"; _f(p3,12,TG)
    arr=sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(5.2),Inches(3.7),Inches(0.8),Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb=AR; arr.line.fill.background()
    rb=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(6.2),Inches(3.0),Inches(5.5),Inches(2.2))
    rb.fill.solid(); rb.fill.fore_color.rgb=LB; rb.line.color.rgb=AR; rb.line.width=Pt(2)
    tf2=rb.text_frame; tf2.word_wrap=True
    p=tf2.paragraphs[0]; p.text="API 实际行为"; _f(p,14,AR,True)
    p2=tf2.add_paragraph(); p2.text="接受 nprobe=0"; _f(p2,13,TD)
    p3=tf2.add_paragraph(); p3.text="→ 搜索执行（错误结果）"; _f(p3,13,AR)
    p4=tf2.add_paragraph(); p4.text="无 crash，无错误码"; _f(p4,12,TG)
    nb=sl.shapes.add_textbox(Inches(0.5),Inches(5.5),Inches(11.5),Inches(0.6))
    np_=nb.text_frame.paragraphs[0]; np_.text="偏离 → 查询语义被破坏；38 个中 37 个不 crash"; _f(np_,14,AR,True)

# === 21 页 ===

s=prs.slides.add_slide(BLANK)
bg=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb=HBG; bg.line.fill.background()
tb=s.shapes.add_textbox(Inches(0.8),Inches(2.4),Inches(11.7),Inches(2.2)); tf=tb.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="TestVDB"; _f(p,54,WH,True)
p2=tf.add_paragraph(); p2.text="用 LLM 发现实际行为偏离文档的缺陷"; _f(p2,22,RGBColor(0xCC,0xDD,0xEE))
m=s.shapes.add_textbox(Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.5))
for i,ln in enumerate(["作者（待定）","单位","会议/Session（待定）"]):
    p=m.text_frame.paragraphs[0] if i==0 else m.text_frame.add_paragraph(); p.text=ln; _f(p,16,RGBColor(0xAA,0xCC,0xDD))

s=slide(2,"开场","VDBMS 缺陷代价高 —— 多数为功能性失败")
image(s,"fig1_rag_arch.png",top=1.3,left=2.5,w=8.3)
table(s,["来源","发现"],[["bug 研究 (Xie 2025)","> 50% 功能性失败"],["测试 roadmap (Wang 2025)","~43% 错误行为；oracle 挑战"],["VDBFuzz (Wang 2026)","首个 fuzzer —— crash oracle"]],top=4.7,h=2.2)

s=slide(3,"问题","文档-实现缺陷：实际行为偏离文档")
bullets(s,[("文档-实现一致性 —— API accept/reject 匹配文档？",TD),("正确性 —— 结果数学正确？",TD)],top=1.3,sz=16)
defect_example(s)

s=slide(4,"问题","38 个中 37 个不崩溃 —— fuzzer 检测不到")
tb=s.shapes.add_textbox(Inches(1),Inches(5.0),Inches(11.3),Inches(1.2))
p=tb.text_frame.paragraphs[0]; p.text="37 / 38"; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(64); p.font.bold=True; p.font.color.rgb=AO
lb=s.shapes.add_textbox(Inches(1),Inches(6.2),Inches(11.3),Inches(0.5))
lp=lb.text_frame.paragraphs[0]; lp.text="确认缺陷不产生 crash —— fuzzer 无法触及"; lp.alignment=PP_ALIGN.CENTER; _f(lp,16,TD)

s=slide(5,"方法","TestVDB：LLM 提取 claim 并裁决一致性"); naive_pipeline(s)
note(s,"crash/差分/蜕变/属性 oracle 都到不了此残差；只有 LLM 能裁决 accept/reject。这个朴素 pipeline 会产生误报 —— LLM judge 不可靠（下页）。")

s=slide(6,"核心洞察","源歧义鸿沟：断言 vs claim"); image(s,"fig3_source_ambiguity_gap.png",top=1.3,left=0.8,w=11.7)

s=slide(7,"核心洞察","家族特定错误：judge 确认 extractor 偏差")
bullets(s,["一个 LLM 家族同时提取 claim 和裁决，共享偏差","judge 确认 extractor 错误 —— self-preference","Panickssery (2024); Wataoka (2024)",("缓解：跨模型验证",AG)])

s=slide(8,"核心洞察","任务内在错误：不同家族推断相同错误 claim"); image(s,"fig4_two_layer_venn.png",top=1.3,left=0.9,w=11.5)

s=slide(9,"核心洞察","跨模型验证覆盖家族特定，不覆盖任务内在")
tc(s,"家族特定",["consistencyLevel —— GLM/DeepSeek 说法不同","跨模型能发现分歧"],
     "任务内在",["timeout —— 都提取 '>= 1'（同错）","跨模型看到一致","只有源码能证伪"],lc=AB,rc=AO)

s=slide(10,"方法","完整 pipeline：源码证伪解决跨模型解决不了的问题")
image(s,"fig5_pipeline.png",top=1.3,left=0.5,w=12.3)
note(s,"加到朴素 pipeline（P6）：dev-reviewer 读源码证伪。12 子句 pilot：跨模型 7/12，源码 12/12（n=29 见 P17）。")

s=slide(11,"方法","证伪规则：shardsNum=0 选默认值")
tc(s,"Over-strict 子句（LLM）",["shardsNum >= 1","probe: shardsNum=0 → API 200","LLM 裁决：'违规'"],
     "源码证伪",["源码: if shardsNum==0 { use default }","0 选默认 —— over-strict","FP 消除","与 MASTOR (as designed) 反向"],lc=AR,rc=AG)
note(s,"+ dev-reviewer 读源码证伪；novelty gate 去重。抑制朴素 pipeline 误报（P6）。")

s=slide(12,"实验","四个 RQ；提交 111，确认 38")
g1=s.shapes.add_table(5,2,Inches(0.4),Inches(1.5),Inches(5.8),Inches(3.8)).table; g1.first_row=True; g1.horz_banding=False
for j,h in enumerate(["RQ","问题"]):
    c=g1.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
    for p in c.text_frame.paragraphs: _f(p,12,WH,True)
for i,(rq,q) in enumerate([["RQ1","多少 doc-impl 缺陷？"],["RQ2","源码是否抑制 FP？"],["RQ3","跨模型 vs 源码？"],["RQ4","无模型子类？"]]):
    for j,v in enumerate([rq,q]):
        c=g1.cell(i+1,j); c.text=v; c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
        for p in c.text_frame.paragraphs: _f(p,11,TD)
g2=s.shapes.add_table(7,3,Inches(6.8),Inches(1.5),Inches(6.1),Inches(4.5)).table; g2.first_row=True; g2.horz_banding=False
for j,h in enumerate(["VDBMS","提交","确认"]):
    c=g2.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=AB
    for p in c.text_frame.paragraphs: _f(p,12,WH,True)
for i,row in enumerate([["Milvus","51","22"],["Weaviate","30","3"],["Qdrant","26","13"],["MeiliSearch","3","0"],["Chroma","1","0"],["合计","111","38"]]):
    for j,v in enumerate(row):
        c=g2.cell(i+1,j); c.text=v; c.fill.solid(); c.fill.fore_color.rgb=LB if i%2==0 else WH
        for p in c.text_frame.paragraphs: _f(p,11,TD)

s=slide(13,"实验","~85% doc-impl 缺陷；VDBFuzz: 0 crash")
tc(s,"组成（非普遍率）",["~85% 文档-实现缺陷","~10% 经典","~5% 并发","确认子集 89%"],
     "VDBFuzz (Qdrant v1.18.2)",["我们跑了 VDBFuzz: 26,000 请求","0 crash, 0 非-200","TestVDB 发现 doc-impl 缺陷","不相交缺陷类"],lc=AB,rc=AG)

s=slide(14,"实验","源码 anchor: 81% FP 抑制（从 31%）")
bn(s,"81%","源码 anchor 抑制的误报","相比 31%；真阳性 96.7%（n=30）")

s=slide(15,"实验","精度随源码 anchor 扩展：25.5% → 45.6% → 69.2%"); image(s,"fig7_ablation.png",top=1.5,left=2.5,w=8.3)

s=slide(16,"实验","RQ3 n=29: 源码 16/16; 显式边界 0/13; κ=1.0")
table(s,["子型","n","任务内在","跨模型","源码"],
      [["参数 over-strict","12","5/12","7/12","12/12"],["行为 over-strict","4","4/4","0/4","4/4"],
       ["显式边界 negative","13","0/13","—","—"],["厂商内对比","—","56% vs 0%","—","—"],
       ["跨模型 κ (n=20)","20","—","κ=1.0","—"]],top=1.5)

s=slide(17,"实验","无模型不变量子类独立发现缺陷")
table(s,["不变量","观测","跨厂商"],[["COSINE 距离","相同向量距离 > 1.0","Milvus+Qdrant"],
     ["索引完整性","索引返回 25 中 2 个","Milvus+Qdrant"],["Payload 过滤","过滤缺失字段返回点","Milvus+Qdrant"]])

s=slide(18,"相关工作","已有工作: 低歧义; TestVDB: 歧义领域")
table(s,["工作线","代表","差异"],
      [["VDBMS 测试","VDBFuzz/roadmap/bug study","crash vs doc-impl"],
       ["REST oracle","AGORA+/SATORI/MASTOR","低歧义；我们用 NL 文档"],
       ["文档 oracle","Toradocu/Doc2OracLL/AugmenTest/Konstantinou","信任 LLM；我们证伪"],
       ["DB 正确性","NoREC/TLP/DQE/DDLCheck","有参考语义；我们没有"],
       ["LLM-judge","Panickssery/Wataoka/Haldar","正交 —— source grounding"]],top=1.4,h=5.6)

s=slide(19,"收尾","威胁与结论")
tc(s,"威胁",["over-strict (n=16) 最 contingent","限 Milvus+Qdrant","机制相关非因果"],
     "结论",["doc-impl 抗确定性检查","LLM: 家族特定+任务内在","超越 VDBMS"],lc=AR,rc=AG)

s=slide(20,"收尾","TestVDB —— 四个核心结果"); image(s,"fig8_summary.png",top=1.4,left=1.5,w=10.3)

prs.save(OUT)
print(f"[wrote] {OUT} ({len(prs.slides._sldIdLst)} slides)")
